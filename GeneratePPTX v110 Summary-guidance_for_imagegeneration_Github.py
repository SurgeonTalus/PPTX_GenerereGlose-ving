import requests  # Import requests library
import json  # Import json library
from pptx import Presentation  # Import Presentation from python-pptx
from pptx.util import Cm
from pptx.dml.color import RGBColor
import os
import base64
from PIL import Image
import numpy as np
from dis_bg_remover import remove_background  # Import dis-bg-remover
import subprocess
import time
import cv2
import traceback
import onnxruntime as ort
from pptx.util import Cm, Pt
import re
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from transformers import pipeline
from PIL import Image, ImageFilter, ImageDraw
import numpy as np
import tkinter as tk
from tkinter import simpledialog
import threading
import time
import wikipedia
from gtts import gTTS

#Draw things must save images to downloads. 
# #Anaconda python 3.11.4 base
#Install anaconda if not working
#press cmd + shift + p, test diffrent python interpreters. Pip install in terminal will revel the system default. Eksample: /opt/anaconda3/lib/python3.12
#pip install wikipedia


#https://docs.drawthings.ai/documentation/documentation/8.scripts/
#https://lmstudio.ai/docs/api/sdk/lmstudioclient

#Steps for image generation. 4 For turbo, 25 for regular. 
STEPS = 4
#Include original slide text when sending to image generator? Works Best for SD3.5. Can be done by specifying in translate prompt.
include_original_text = False
REMOVE_BACKGROUND = True  # Toggle for background removal

# Summary uses gemma-3-12b-it
# Define available models
available_models = [
    {"id": "gemma-3-12b-it", "name": "gemma-3-12b-it"}, #Run it with   Exprimental Kash q4 to reduce memory from 10gb to 7gb.
    {"id": "qwen2.5-14b-instruct-mlx", "name": "Qwen 2.5 14B Instruct MLX"},
    {"id": "mistral-nemo-instruct-2407", "name": "mistral-nemo-instruct-2407"},
    {"id": "mistral-small-24b-instruct-2501", "name": "mistral-small-24b-instruct-2501"},
    {"id": "qwen2.5-7b-instruct-1m", "name": "qwen2.5-7b-instruct-1m"},
    {"id": "qwen2.5-7b-instruct-1m", "name": "If Already Loaded Into Memory it Will Overwrite this choise, press whatever"}
]


# Display model options
print("Select a model to load: (Add models manually from line 35) Qwen 2.5 14B Instruct MLX is superoir to Mistral but not perfect."
      "Turn on KV cache Quantization = 4000, k6 on 16GB systems for the models to avoid limiting Context Window")
for index, model in enumerate(available_models, start=1):
    print(f"{index}. {model['name']}")

# Get user selection
selected_index = input("Enter the number of the model to load: (Already loaded models will overwrite)")

def unload_models(unload_all=False):  # DONT DO CHANGES HERE, modify "unload_models(unload_all=False)" above"
    try:
        if unload_all:  # If 'unload_all' is True, proceed with unloading all models | offload
            # Run the command "lms unload --all" to unload all models 
            result = subprocess.run(["lms", "unload", "--all"], check=True, capture_output=True, text=True)
            print(result.stdout)  # Print the output of the command execution
        else:
            print("unload_all=False VERY RAM INTENSIVE! TO AVOID CRASHES, SET TO TRUE")  # If 'unload_all' is False, do nothing and notify the user
    except subprocess.CalledProcessError as e:  # Handle errors if the command execution fails
        print(f"Error unloading model(s): {e.stderr}")  # Print the error message  

def select_temperature():
    options = {
        1: (0.1, "For gemma-3-12b-it. Likely Repetitive for others (Combine with adjusting repetition penalty manually in LM Studio)"),
        2: (0.15, "Optimal for Mistral Small 24b, Spesielt på Norsk."),
        3: (0.3, "Mistral Nemo Recommended setting"),
        4: (0.4, ""),
        5: (0.5, ""),
        6: (0.6, ""),
        7: (0.7, "Default Norm"),
        8: (0.85, "Creative"),
        9: (1, "Too creative, likely Gibberish grammar"),
        10: (2, "Brace for some interesting language mixing"),
    }

    print("Select a temperature setting:")
    for key, (value, description) in options.items():
        print(f"{key}: Temperature = {value} {description}")

    choice = int(input("\nEnter the number of your choice: "))

    temperature = options.get(choice, (0.7, "Default Norm"))[0]
    print(f"\nSelected Temperature for the Text Prosessor. For imagegenerator, define this manually in code. Nemo 0.7 by default: {temperature}")

    return temperature


if __name__ == "__main__":
    temperature = select_temperature()
    
# Path to the Draw Things app
draw_things_app_path = "/Applications/Draw Things.app"
print("Remember to start API server inside of Draw Things app")
        
# Open Draw Things app
def open_draw_things():
    try:
        subprocess.Popen(["open", draw_things_app_path])
        print("Launching Draw Things app...")
        # Wait a few seconds to allow the app to fully initialize
        time.sleep(0.1)
        print("turn on HTTP manually as it disables on every start")
    except Exception as e:
        print(f"Failed to open Draw Things app: {e}")

        
# Validate input and get the model ID
try:
    selected_model = available_models[int(selected_index) - 1]["id"]
except (IndexError, ValueError):
    print("Invalid selection.")
    selected_model = None

# Define the global variable at the top of your script
twoLinesSummaryOfTheWholeTextForLMstudio = ""

# Function to communicate with the local API
def send_request_to_localhost(data):
    if not selected_model:
        print("No valid model selected.")
        return None

    # Check available models
    models_response = requests.get("http://localhost:1234/api/v0/models")
    if models_response.status_code == 200:
        models = models_response.json().get("data", [])
        loaded_model = next((model for model in models if model["state"] == "loaded"), None)
        
        # If no model is loaded, try to load the selected model
        if not loaded_model:
            print(f"No loaded model found. Attempting to load {selected_model}.")
            model_to_load = next((model for model in models if model["id"] == selected_model), None)
            if model_to_load:
                load_response = requests.post(
                    "http://localhost:1234/api/v0/model/load",
                    json={"model": model_to_load["id"]}
                )
                if load_response.status_code == 200:
                    print(f"Model {selected_model} loaded successfully.")
                    current_model = model_to_load["id"]
                else:
                    print(f"Failed to load model with status code: {load_response.status_code}")
                    return None
            else:
                print(f"Model {selected_model} not found.")
                return None
        else:
            current_model = loaded_model["id"]
    else:
        print(f"Failed to list models with status code: {models_response.status_code}")
        return None

    url = "http://localhost:1234/v1/chat/completions"
    request_data = {
        "model": current_model,
        "messages": data["messages"],
        "temperature": temperature,
        "max_tokens": 20000,
        "stream": True,
        "reset_history": True,
        "noHup": False
    }
    




    response_text = ""

    # Function to process the response stream
    def process_stream(response):
        nonlocal response_text
        for chunk in response.iter_content(chunk_size=None):
            if chunk:
                chunk_data = chunk.decode('utf-8').strip()
                if chunk_data.startswith("data: "):
                    chunk_data = chunk_data[6:]
                try:
                    chunk_json = json.loads(chunk_data)
                    if "choices" in chunk_json:
                        for choice in chunk_json["choices"]:
                            if "delta" in choice and "content" in choice["delta"]:
                                response_text += choice["delta"]["content"]
                except json.JSONDecodeError:
                    response_text += chunk_data
                print("Partial response:", response_text)

        # After processing, save response_text to a .txt file in the script's directory
        save_response_to_file(response_text)

    # Function to save response_text to a .txt file
    def save_response_to_file(response_text):
        # Get the current script's directory
        script_dir = os.path.dirname(os.path.realpath(__file__))
        file_path = os.path.join(script_dir, "response.txt")
        
        # Write the response text to a file
        with open(file_path, "w") as file:
            file.write(response_text) #Saves the whole response to a txt file stored in the variable response_text
        
        print(f"Response saved to: {file_path}")


    # First request pass
    print("Starting first request...")
    response = requests.post(url, json=request_data, stream=True)
    if response.status_code == 200:
        try:
            process_stream(response)
        except requests.exceptions.RequestException as e:
            print(f"Request exception: {e}")

    # Check for the number of '#' characters
    hashtag_count = response_text.count('#')
    print(f"Number of '#' characters: {hashtag_count}")

    # If two or fewer '#' characters, make a second request with additional instruction
    if hashtag_count <= 2:
        print("Less than or equal to 2 '#' found. Starting second pass with additional instruction...")

        # Add the additional instruction to the conversation history
        additional_instruction = {
            "role": "user",
            "content": "Divide text into smaller sections with the # hashtag symbol."
        }
        new_messages = data["messages"] + [additional_instruction]

        second_request_data = {
            "model": current_model,
            "messages": new_messages,
            "temperature": temperature,
            "max_tokens": 20000,
            "stream": True,
            "reset_history": False,  # Keep conversation context
            "noHup": False
        }

        second_response = requests.post(url, json=second_request_data, stream=True)
        if second_response.status_code == 200:
            try:
                process_stream(second_response)
            except requests.exceptions.RequestException as e:
                print(f"Request exception on second pass: {e}")

    print("Final response:", response_text)
    
    ############SUMMARY#####################
    #####Load txt file as variable


    script_dir = os.path.dirname(os.path.abspath(__file__))  # Get script's directory
    file_path = os.path.join(script_dir, "response.txt")  # Construct full path

    with open(file_path, "r", encoding="utf-8") as file:
        globalresponse_text = file.read()

    print(globalresponse_text)  # Verify content
    
    time.sleep(2)
    # Attempt to load the file
    globalresponse_text = ""

    if os.path.exists(file_path):
        try:
            with open(file_path, "r", encoding="utf-8") as file:
                globalresponse_text = file.read()
        except Exception as e:
            print(f"Error reading response.txt: {e}")
    else:
        print("response.txt not found.")

    print(globalresponse_text)  # Verify content

    
    




    # Function to get a summary description from LM Studio
    def description_prompt_text(globalresponse_text):
        if not globalresponse_text.strip():
            print("No text available to summarize.")
            return "No description generated."

        request_data = {
            "messages": [
                {"role": "system", "content": "Do not comment. Provide a descriptive one-sentence summary in english. Then one sentence with 'the purpose is likely to illustrate'"},
                {"role": "user", "content": globalresponse_text}
            ],
            "model": "gemma-3-12b-it",
            "temperature": 0.1,
            "max_tokens": 500,
            "stream": False,
            "reset_history": True,
        }
        
        url = "http://localhost:1234/v1/chat/completions"  # Change this to your LM Studio API URL
        headers = {"Content-Type": "application/json"}

        try:
            response = requests.post(url, json=request_data, headers=headers)
            response.raise_for_status()  # Raise an error for bad responses (4xx, 5xx)
            response_json = response.json()
            return response_json.get("choices", [{}])[0].get("message", {}).get("content", "No description generated.")
        except requests.exceptions.RequestException as e:
            print(f"Error generating description: {e}")
            return "No description generated."



    global twoLinesSummaryOfTheWholeTextForLMstudio  # Declare as global to modify it
    twoLinesSummaryOfTheWholeTextForLMstudio = description_prompt_text(globalresponse_text)

    print("############# SUMMARY ##################" + "\n" + 
        "If gemma-3-12b-it is installed:" + "\n" + 
        "Description going to LM studio to summarize the whole script: " + "\n" + 
        twoLinesSummaryOfTheWholeTextForLMstudio + "\n" + 
        "###############################" + "\n")

    ########SUMMARYEND#############
  
    return response_text


# Function to generate image using Draw Things API
DRAW_THINGS_URL = 'http://127.0.0.1:7860/sdapi/v1/txt2img'
#IMG_SIZE = 512
  # Define steps outside the function

def generate_image(prompt):
    print(f"Generating image with {STEPS} steps...")  # Print the number of steps
    params = {
        "prompt": prompt,
        "negative_prompt": "(bokeh, worst quality, low quality, normal quality, (variations):1.4), blur:1.5",
        "seed": 4068245935,
        "steps": STEPS,  # Use the external variable
        "guidance_scale": 10, # Text guidance, guide, 4.5
        "batch_count": 1
    }


    headers = {"Content-Type": "application/json"}
    response = requests.post(DRAW_THINGS_URL, json=params, headers=headers)
    if response.status_code == 200:
        data = response.json()
        images = data.get("images", [])
        if images:
            temp_image_path = os.path.join("/tmp", "generated_image.png")
            with open(temp_image_path, "wb") as img_file:
                img_file.write(base64.b64decode(images[0]))

            # Attempt to remove background
            if REMOVE_BACKGROUND:
                processed_image_path = remove_background_from_image(temp_image_path)
                return (temp_image_path, processed_image_path)  # Return both paths
            
            else:
                print("Background removal disabled, using the original image.")
            
            return temp_image_path, None  # Return original and none for background removed

    else:
        print(f"Error generating image: {response.status_code}, {response.text}")

    return None, None  # Return both as None if image generation failed

def normalize(image, mean, std):
    """Normalize a numpy image with mean and standard deviation."""
    return (image / 255.0 - mean) / std

# Function to remove background from image
def remove_background_from_image(image_path):
    model_path = os.path.expanduser("~/Downloads/isnet_dis.onnx")
    # If your script is in a directory relative to the ONNX file, you can use: model_path = os.path.join(os.getcwd(), "isnet_dis.onnx")
    print("Download isnet_dis.onnx to the downloads folder in case of error")
    try:
        extracted_img, mask = remove_background(model_path, image_path)
        
        if extracted_img is None:
            print("Failed to remove background: No output image")
            return None

        extracted_img_pil = Image.fromarray(extracted_img)
        
        # Get the Downloads folder path and create a filename
        downloads_folder = os.path.expanduser("~/Downloads")
        base_filename = "processed_image_no_bg.png"
        processed_image_path = os.path.join(downloads_folder, base_filename)

        extracted_img_pil.save(processed_image_path)
        print(f"Background removed and image saved to {processed_image_path}")
        return processed_image_path
    except Exception as e:
        print(f"Failed to remove background: {e}")
        return None

import os
import cv2
import numpy as np
import onnxruntime as ort
from PIL import Image
import traceback

# Updated remove_background function
def remove_background(model_path, image_path):
    if model_path is None or image_path is None:
        return None, None

    input_size = (1024, 1024)

    try:
        # Load the ONNX model
        session = ort.InferenceSession(model_path)
        im = cv2.imread(image_path, cv2.IMREAD_COLOR)
    #    im = cv2.cvtColor(im, cv2.COLOR_BGR2RGB)  # Convert from BGR to RGB if using OpenCV

        # If image is grayscale, convert to RGB
        if len(im.shape) == 2:
            im = cv2.cvtColor(im, cv2.COLOR_GRAY2RGB)
        
        # Print initial image shape and values
        print(f"Original Image shape: {im.shape}")
        
        # Normalize the image using NumPy
        im = im.astype(np.float32)  # Convert to float
        im_normalized = normalize(im, mean=[0.5, 0.5, 0.5], std=[1.0, 1.0, 1.0])
            
        # Resize the image
        im_resized = cv2.resize(im_normalized, input_size, interpolation=cv2.INTER_LINEAR)
        im_resized = np.transpose(im_resized, (2, 0, 1))  # CHW format
        im_resized = np.expand_dims(im_resized, axis=0)  # Add batch dimension

        # Run inference
        im_resized = im_resized.astype(np.float32)  
        ort_inputs = {session.get_inputs()[0].name: im_resized}
        ort_outs = session.run(None, ort_inputs)
            
        # Process the model output
        result = ort_outs[0][0]  # Assuming single output and single batch
        result = np.clip(result, 0, 1)  # Assuming you want to clip the result to [0.5, 1] is half transparent
        result = (result * 255).astype(np.uint8)  # Rescale to [0, 255]
        result = np.transpose(result, (1, 2, 0))  # HWC format
        
        # Debug: Check result image
        print(f"Result image shape after background removal: {result.shape}")
        
        # Resize to original shape
        original_shape = im.shape[:2]
        result = cv2.resize(result, (original_shape[1], original_shape[0]), interpolation=cv2.INTER_LINEAR)

        # Ensure 'result' is 2D (H x W) and add an axis to make it (H x W x 1)
        alpha_channel = result[:, :, np.newaxis]

        # Debug: Check alpha channel
        print(f"Alpha channel shape: {alpha_channel.shape}")
        
        # Concatenate the RGB channels of 'im' with the alpha channel
        im_rgba = np.concatenate((im, alpha_channel), axis=2)

        # Debug: Check RGBA image before conversion
        print(f"RGBA Image shape: {im_rgba.shape}")
        
        # Convert to BGRA (Blue, Green, Red, Alpha)
        im_bgra = cv2.cvtColor(im_rgba, cv2.COLOR_RGBA2BGRA)

        # Debug: Check BGRA image
        print(f"BGRA Image shape: {im_bgra.shape}")
        
        # Convert to uint8 before returning
        im_bgra = im_bgra.astype(np.uint8)

        return im_bgra, result
    except Exception as e:
        print("An error occurred:")
        traceback.print_exc()
        return None, None

# Function to calculate average color
def calculate_average_color(image_path):
    try:
        image = Image.open(image_path)
        image = image.convert("RGB")
        pixels = np.array(image)
        avg_color = np.mean(pixels, axis=(0, 1))
        return tuple(avg_color.astype(int))
    except Exception as e:
        print(f"Error calculating average color: {e}")
        return (255, 255, 255)  # Default to white if there's an error

# Function to brighten a color
def brighten_color(color, factor=2):
    r, g, b = color
    r = min(int(r * factor), 255)
    g = min(int(g * factor), 255)
    b = min(int(b * factor), 255)
    return r, g, b

# Function to set background color for slides
def set_background_color(slide, image_path):
    try:
        avg_color = calculate_average_color(image_path)
        brightened_color = brighten_color(avg_color, factor=3)

        # Apply the brightened color as the slide's background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*brightened_color)
        print(f"Set background color: {brightened_color}")
    except Exception as e:
        print(f"Failed to set background color: {e}")



def apply_semi_oval_blur(image_path):
    # Open the original image
    img = Image.open(image_path)
    width, height = img.size

    # Create a new image with transparent background (RGBA mode)
    img_with_alpha = img.convert("RGBA")
    
    # Create a rounded rectangle gradient mask
    mask = Image.new("L", (width, height), 255)  # Start with a white mask (full opacity)
    draw = ImageDraw.Draw(mask)
    
    # Define corner radius (adjust as needed)
    corner_radius = min(width, height) // 2
    
    # Draw a filled rounded rectangle with full opacity
    draw.rounded_rectangle([(0, 0), (width, height)], corner_radius, fill=0)
    
    # Apply a gradient effect from the edges to the center
    for y in range(height):
        for x in range(width):
            # Calculate the distance from the nearest edge
            edge_dist_x = min(x, width - x)
            edge_dist_y = min(y, height - y)
            edge_dist = min(edge_dist_x, edge_dist_y)
            
            # Calculate the transparency based on distance from the edges
            alpha = max(0, min(255, int((edge_dist / corner_radius) * 255)))
            mask.putpixel((x, y), alpha)
    
    # Apply the mask to blend with alpha transparency near the edges
    img_with_alpha.putalpha(mask)  # Apply the transparency mask

    # Save the modified image
    downloads_folder = os.path.expanduser("~/Downloads")
    blurred_image_path = os.path.join(downloads_folder, "feather_falloff_image.png")

    img_with_alpha.save(blurred_image_path)  # Save the modified image

    print(f"Image with feather falloff saved to: {blurred_image_path}")

    return blurred_image_path  # Return the path to the modified image


def remove_bulletpoints_after_exclamation(prs):
    # Iterate over all slides in the presentation
    for slide in prs.slides:
        # Iterate over all shapes in the slide
        for shape in slide.shapes:
            # Check if the shape has text (e.g., a text box or placeholder)
            if hasattr(shape, "text"):
                # Split the text at the first occurrence of "!!!"
                if "!!!" in shape.text:
                    shape.text = shape.text.split("!!!")[0]
        # Function to brighten a color
                                 
def create_pptx(response_text):
    prs = Presentation()
    # Set the slide width and height for a 16:9 aspect ratio
    prs.slide_width = Cm(33.867)  # 16:9 width
    prs.slide_height = Cm(19.05)  # 16:9 height
    slide_layout = prs.slide_layouts[1]  # Using a predefined layout with title and subtitle placeholders

#!!!manuallyLoadedFirstSlide_text = ("""#""" + response_text) loads at the start of the first slide essesial to force creation of the first powerpoint slide and avoid subtitle_placeholder = slide.placeholders[1] errir
#!!!manuallyLoadedFirstSlide_text = ("""#""" + response_text) loads at the start of the first slide essesial to force creation of the first powerpoint slide.and avoid subtitle_placeholder = slide.placeholders[1] errir
#!!!manuallyLoadedFirstSlide_text = ("""#""" + response_text) loads at the start of the first slide essesial to force creation of the first powerpoint slide.and avoid subtitle_placeholder = slide.placeholders[1] errir
    manuallyLoadedFirstSlide_text = ("""#""" + response_text + ""
                                     "\n# Slutt""") ##Use for triggering tagging consisten Characters.  Will appear as heading for first slide.
    input_lines = manuallyLoadedFirstSlide_text.splitlines()  ###### Laste inn lokal tekst variabel!
    
    title = ""
    subtitle = ""
    skip_description = "!!!"
    
        # Define save location
    downloads_folder = os.path.expanduser("~/Downloads")
    base_filename = "FirstTestPageScriptWorking_PPTX"
    pptx_filename = os.path.join(downloads_folder, f"{base_filename}.pptx")
    
    # Ensure a unique filename if one already exists. THis does not affect counting in filenames when deleted fore some reason. 
    counter = 1
    while os.path.exists(pptx_filename):
        pptx_filename = os.path.join(downloads_folder, f"{base_filename}_{counter}.pptx")
        counter += 1

 # Error handling - save progress if crash occurs
    try:
        for line in input_lines:
            if line.startswith(('#', '##', '###')):
                # Save previous slide if it exists
                if subtitle:
                    subtitle_placeholder = slide.placeholders[1]
                    subtitle_placeholder.text = subtitle
                    adjust_subtitle_font(subtitle_placeholder)
                    subtitle = ""

                # Create new slide
                slide = prs.slides.add_slide(slide_layout)
                title = line.lstrip('#').strip()
                slide.shapes.title.text = title

                # Save after each slide
                prs.save(pptx_filename)
                print(f"Auto-saved after slide {len(prs.slides)}")

            else:
                subtitle += line + "\n"

        # Save last slide
        if subtitle:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle
            adjust_subtitle_font(subtitle_placeholder)

        # Final save
        prs.save(pptx_filename)
        print(f"Presentation successfully saved as {pptx_filename}")

    except Exception as e:
        print(f"Error occurred: {e}")
        prs.save(pptx_filename)  # Save progress before crash
        print(f"Saved progress before crash: {pptx_filename}")

    for line in input_lines:
        if (line.startswith("***")) or \
        (line.startswith("**")) or \
        (line.startswith("*") and line.endswith("*")):
            line = "##" + line.strip("*")  # Remove leading and trailing * and add ##
        print(line)
    
        if line.startswith(('#', '##', '###')):
            # If a new title is found, create a new slide
            if subtitle:
                print ( "AccurateSubtitle", subtitle)
                subtitleForTTS = title + ". " + subtitle
                subtitle_placeholder = slide.placeholders[1] 

                subtitle_placeholder.text = subtitle

                # Adjust dimensions and position for the last slide's subtitle
                subtitle_placeholder.width = Cm(18)  # Width adjustment
                subtitle_placeholder.height = Cm(15)  # Height adjustment
                subtitle_placeholder.left = Cm(0)  # Left position
                subtitle_placeholder.top = Cm(3.5)  # Top position

                # Set the font size of the subtitle
                text_frame = subtitle_placeholder.text_frame
                initial_font_size = 32
                max_font_size = initial_font_size
                min_font_size = 18  # Minimum font size for scaling

                subtitle_text = "".join([paragraph.text for paragraph in text_frame.paragraphs])
                text_length = len(subtitle_text)

                if text_length > 200:  # Adjust this threshold as needed
                    max_font_size = min(max_font_size - (text_length // 50), initial_font_size)

                max_font_size = max(max_font_size, min_font_size)

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(max_font_size)

###############TTS################
                # Extract the plain text without any formatting for TTS
                

                # Define the file paths in the Downloads folder
                downloads_folder = os.path.expanduser("~/Downloads")
                text_file_path = os.path.join(downloads_folder, 'plain_text.txt')
                audio_file_path = os.path.join(downloads_folder, 'audio_output.aiff')

                # Save the plain text to a .txt file
                with open(text_file_path, 'w', encoding='utf-8') as text_file:
                    text_file.write(subtitleForTTS)

                # Use macOS 'say' command to generate an audio file from the text file
                os.system(f"say -v Henrik -f {text_file_path} -o {audio_file_path}")
        


                # Add the audio file (embedding audio)
                audio_shape = slide.shapes.add_movie(
                    audio_file_path,
                    left=Inches(4.17),
                    top=Inches(6.74),
                    width=Inches(1.67),
                    height=Inches(0.76),
                    poster_frame_image=None,
                    mime_type='audio/aiff'  # Adjust mime type based on the audio file format
                )
################ v95 Bold, Italic, and Underscore: Markdown Conversion ################
   
# Extract the text from paragraphs and generate an audio file
                for paragraph in text_frame.paragraphs:
                    # Get the full text and split it based on **bold**, *italic*, and _underscore_ syntax
                    full_text = paragraph.text
                    paragraph.clear()  # Clear existing text in the paragraph

                    # Split the full_text by **bold**, *italic*, and _underscore_, preserving the formatting
                    segments = re.split(r"(\*\*.*?\*\*|\*.*?\*|_.*?_)", full_text)  # Split and keep bold, italic, and underscore segments

                    for segment in segments:
                        if segment.startswith("**") and segment.endswith("**"):
                            # Apply bold formatting to text inside **
                            run = paragraph.add_run()
                            run.text = segment.strip("*")  # Remove the ** around the text
                            run.font.bold = True  # Apply bold formatting
                        elif segment.startswith("*") and segment.endswith("*"):
                            # Apply italic formatting to text inside *
                            run = paragraph.add_run()
                            run.text = segment.strip("*")  # Remove the * around the text
                            run.font.italic = True  # Apply italic formatting
                        elif segment.startswith("_") and segment.endswith("_"):
                            # Apply underscore formatting to text inside _
                            run = paragraph.add_run()
                            run.text = segment.strip("_")  # Remove the _ around the text
                            run.font.underline = True  # Apply underscore formatting
                        else:
                            # Add normal text (no formatting)
                            run = paragraph.add_run()
                            run.text = segment

                        run.font.size = Pt(max_font_size)



################ Table Detection and Conversion ################ Not working. Migt be bullet points causing error. Test otuside this script. 

                    # Regular expression to detect markdown tables
                    table_regex = r"(\|.*?\|[\r\n]*)+"

                    # Collect text from all paragraphs and reconstruct potential tables
                    full_text = "\n".join([p.text for p in text_frame.paragraphs])  # Join paragraphs to ensure full capture

                    # Check if the full text contains a markdown-style table
                    table_match = re.search(table_regex, full_text)
                    if table_match:
                        table_text = table_match.group()  # Extract matched table text

                        # Split the markdown table into rows
                        rows = table_text.strip().split("\n")

                        # Ensure valid table format (at least header + one row)
                        if len(rows) < 2:
                            raise ValueError("Invalid Markdown table format. Must have at least a header and one data row.")

                        # Extract header row and data rows
                        header_row = rows[0]
                        separator_row = rows[1]  # The `|---|---|` row (ignored)
                        data_rows = rows[2:]  # Actual table data

                        # Determine number of columns correctly
                        num_cols = len([col for col in header_row.split('|') if col.strip()])  # Ignore empty splits
                        num_rows = len(data_rows) + 1  # Header + data rows

                        # Define table position and size in PowerPoint
                        x, y, cx, cy = Inches(1), Inches(1), Inches(6), Inches(3)

                        # Add table to the slide
                        table_shape = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy)
                        table = table_shape.table  # Get the actual table object

                        # Process the header row (bold formatting)
                        header_cells = [col.strip() for col in header_row.split('|') if col.strip()]  # Remove empty splits
                        for j, col in enumerate(header_cells):
                            cell = table.cell(0, j)
                            cell.text = col
                            cell.text_frame.paragraphs[0].font.bold = True  # Apply bold formatting for headers

                        # Process each data row
                        for i, row in enumerate(data_rows, 1):  # Start at 1 (since 0 is header)
                            columns = [col.strip() for col in row.split('|') if col.strip()]
                            for j, col in enumerate(columns):
                                table.cell(i, j).text = col  # Set text in table cell

      


                # Reset subtitle for the next title/ next page
                subtitle = ""

                
            # Create a new slide with the layout for title + subtitle
            prs.slide_width = Cm(33.87)  # 16:9 width
            prs.slide_height = Cm(19.05)  # 16:9 height
            slide = prs.slides.add_slide(slide_layout)
            title = line.lstrip('#').strip()  # Extract the title

            # Extract the title (title)
            title = line.lstrip('#').strip()  # Extract the title
            hyperlink = f'https://www.google.com/search?q={title}'

            # Add title with "->" symbol containing the hyperlink

            title_shape = slide.shapes.title
            if title_shape:
                # Set the title box dimensions
                title_shape.text = title  # Set the title text
                
                # Calculate the font size based on the length of the title
                title_length = len(title)
                if title_length <= 20:
                    font_size = 44
                else:
                    # Reduce font size by 1 for every 2 characters beyond 20, with a minimum of 26
                    extra_chars = title_length - 20
                    font_size = max(26, 44 - (extra_chars // 2))  # Reduce for every other character beyond 20
                    
                
                # Apply the font size to the title
                run = title_shape.text_frame.paragraphs[0].runs[0]  # Access the first run
                run.font.size = Pt(font_size)
                run.font.bold = True  # Set the font to bold
                
                
                #title_shape.width = Cm(21)  # Width: 21 cm
                #title_shape.height = Cm(1.5)  # Height: 3.5 cm
                
                # Add "->" symbol to the same line as the title
                run = title_shape.text_frame.paragraphs[0].add_run()
                run.text = " ->"  # Append the "->" symbol
                run.font.size = Pt(30)
                run.font.underline = False  # Remove underline
                
                # Apply the hyperlink to the "->" part
                title_shape.text_frame.paragraphs[0].runs[-1].hyperlink.address = hyperlink

            # Adds Wikipedia to comments for fact-checks
            def search_wikipedia(title, get_full_text=False): #Ikke endre her, endre i bunn av def med get_full_text = False
                try:
                    # Retry logic for up to 5 seconds if internet is not working
                    start_time = time.time()
                    while time.time() - start_time < 1:
                        try:
                            # Search Wikipedia
                            search_results = wikipedia.search(title)
                            
                            if not search_results:
                                print("No results found.")
                                return None
                            
                            # Get the most relevant result (first result)
                            page_title = search_results[0]
                            page_url = wikipedia.page(page_title).url  # Get the URL of the page
                            
                            if get_full_text:
                                # Get the full page content
                                page = wikipedia.page(page_title)
                                print(f"Title: {page.title}\n")
                                print(page.content)
                                slide.notes_slide.notes_text_frame.text = page.content
                            else:
                                # Get the page summary
                                summary = wikipedia.summary(page_title)
                                print(f"Title: {page_title}\n")
                                print(summary)
                                slide.notes_slide.notes_text_frame.text = summary
                            
                            # Append the Wikipedia URL to slide notes
                            slide.notes_slide.notes_text_frame.text += f"\n\nFor more details, visit: {page_url}"
                            
                            return  # Successfully completed, exit the loop
                            
                        except wikipedia.exceptions.DisambiguationError as e:
                            print("Disambiguation Error: Multiple possible results. Try being more specific.")
                            print(e.options)
                            return None
                        except wikipedia.exceptions.PageError:
                            print("Page not found.")
                            return None
                        except Exception as e:
                            print(f"An error occurred: {e}")
                            time.sleep(1)  # Wait before retrying
                            
                    print("Failed to fetch Wikipedia content within 5 seconds.")
                
                except Exception as e:
                    print(f"An unexpected error occurred: {e}")
                    return None        
        

            # Set the Wikipedia language to Norwegian (Bokmål) if needed
            wikipedia.set_lang("no")

            # Specify whether you want the full text or the summary (True for full text, False for summary)
            get_full_text = False  # Change this variable to True for full text, False for summary

            # Search for the given title
            search_wikipedia(title, get_full_text)
            
            # slide.notes_slide.notes_text_frame.text = "https://www.google.com/search?q=" + title


        else:
            # If it's not a title, accumulate it as a subtitle
            subtitle += line + "\n"

    # Add the last subtitle to the final slide after the loop
    if subtitle:
        subtitle_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None ##548 Skip error?
        subtitle_placeholder.text = subtitle

        # Adjust dimensions and position for the last slide's subtitle
        subtitle_placeholder.width = Cm(18)  # Width adjustment
        subtitle_placeholder.height = Cm(15)  # Height adjustment
        subtitle_placeholder.left = Cm(0)  # Left position
        subtitle_placeholder.top = Cm(4)  # Top position

         ##Set the font size of the subtitle
        text_frame = subtitle_placeholder.text_frame
        initial_font_size = 32
        max_font_size = initial_font_size
        min_font_size = 18  # Minimum font size for scaling

        subtitle_text = "".join([paragraph.text for paragraph in text_frame.paragraphs])
        text_length = len(subtitle_text)

        if text_length > 100:  # Adjust this threshold as needed
            max_font_size = min(max_font_size - (text_length // 50), initial_font_size)

        max_font_size = max(max_font_size, min_font_size)
 
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(max_font_size)
                
 
    
    # Now generate image for each slide based on the content (title + subtitle)
    
    

    for i, slide in enumerate(prs.slides):
        subtitle = slide.shapes.placeholders[1].text  # Get the subtitle from the placeholder
        prompt_text = "Cartoon for tweens." + slide.shapes.title.text + "\n" + subtitle  # Combine title and subtitle

#####Add extra instructions for Draw Things. #######################
#####Add extra instructions for Translation.#######################
#####Add extra instructions for Scene..#######################
#####Add extra instructions for LLM..#######################
#####Translation

# Store the time of the last loop run (initialize it with the current time)
    last_run_time = time.time()
    ##### Deactivate by setting # before + like this: #+
    IMGinstruct = " " + "Here comes instructions: Based on the given input text write out details for one illustration. Sometimes the input text sounds like an instruction, in that case, illustrate the instruction."
    IMGsummary = " " + "Background context for the input, this is supposed to guide the image description in cases where input does not make it clear:" + twoLinesSummaryOfTheWholeTextForLMstudio
    IMGlanguage = " " + "answer in english"
    IMGEducatioanl = " " + "Focus on illustrating the Verbs or adjectives describing the noun, not the main noun itself if any."
    IMGnoChatting = " " + "You are given instructions to execute,  no comments, just prossesed description/results."
    IMGautonomy = " " + "Deside witch character or object is the most important to illustrate based on the provided text. Be very explicit of your choise!"
    IMGpreferObjects = " " #+ "Objects are strongly preferd over humans"
    IMGifCharacters = " " + "If characters, Do not describe mutlible characters, skip all metions of other characters afeter you have chosen. If describing objects, never describe who owns it. "
    IMGifHumans = " " + "If human, describe ""facial expression"" in one word (Happy, sad, scared etc). If human, consider to include ""feets on ground"" as this will create a full body view instead of default portrait."
    IMGmood = " " + "Describe the mood of the scene in one word."
    IMGobjects = " " + "Make shure to describe objects that are important for the story in a way that makes them very noticable/apearent for the eye, those things should be big."
    IMGconsise = " " + "Write everything in a few short consise lines . ""Never write unessesary things like - **Description:** X not visible in this illustration."" as this is  meaningless fill information."
    IMGbacground = " " #+ "Give simple descriptions in complete sentences of type of scene, eksample; kitchen, store, office, woods, cabin etc, then add some few words of complementary details."
    IMGincludeFilenameToForceTopic = " " #+ base_filename
    #DISABLE A VARIABLE
    IMGexample = " " #Comment out the rest
    
    def translate_prompt_text(prompt_text):
        # Send the prompt text to LM Studio for translation
        url = "http://localhost:1234/v1/chat/completions"
        request_data = {
            "messages": [{"role": "system", "content": IMGincludeFilenameToForceTopic + IMGinstruct + IMGsummary + IMGnoChatting + IMGlanguage + IMGEducatioanl + IMGautonomy + IMGpreferObjects + IMGifCharacters + IMGifHumans + IMGmood + IMGobjects + IMGconsise + IMGbacground + ""}],
            "model": "gemma-3-12b-it",  # Assuming the model id
            "temperature": 0.1,
            "max_tokens": 20000,
            "stream": False,
            "reset_history": True,
        }
        request_data["messages"].append({"role": "user", "content": prompt_text })
        
        # Send request to LM Studio for translation
        response = requests.post(url, json=request_data)
        if response.status_code == 200:
            response_json = response.json()
            translated_text = response_json.get('choices', [{}])[0].get('message', {}).get('content', '')
            # Unload Text Model before loading Image Model
            return translated_text
        else:
            print(f"Translation failed. Status code: {response.status_code}")
            return prompt_text  # Return the original prompt text if translation fails
    unprossessed_prompt_text = ""
    
    # Iterate through the slides to prepare and send the requests
    for i, slide in enumerate(prs.slides):
        # Track the time since the last loop run
        current_time = time.time()
        time_diff = current_time - last_run_time
        print(f"Generation Time: {time_diff:.2f} seconds")
        print(f"Generation Time: {time_diff:.2f} seconds")
        print(f"Generation Time: {time_diff:.2f} seconds")

        # Update last run time for the next iteration
        last_run_time = current_time

        subtitle = slide.shapes.placeholders[1].text  # Get the subtitle from the placeholder
        prompt_text = "" + slide.shapes.title.text + "\n" + subtitle  # Combine title and subtitle
        unprossessed_prompt_text = "" + slide.shapes.title.text + "\n" + subtitle #Will be added to image generator if specifyed
        print ("Raw unprossesed slides now prosessing" + unprossessed_prompt_text)
        # Translate the prompt_text before sending it
            
    
        translated_prompt_text = translate_prompt_text(prompt_text)
        
        # Unload Text Model before loading Image Model
        unload_models(unload_all=False)
        print("")
        print("SLOWER, REDUCE MEMORY, AVOID CRASH (Unload LLM before image generation), SET: unload_all=True")
        print("")
        print("FASTER, HAVE LLM AND IMAGE GENERATOR LOADED IN PARRALELL, SET: unload_all=False")
        print("")
        print(f"TURBO MUST BE 4! Currently {STEPS} steps...")  # Print the number of steps



################# Prompt Directly into DrawThings (dt) ################
        dtEksample = "" #+ "set to blank to avoid." + " "
        dtFilename = "" #+ base_filename
        dtStyle = ""#+ "" #Fill with style from here:  style prompts here: https://rikkar69.github.io/SDXL-artist-study/
        dtDisney80 = "" #+ "Don Bluth Disney 80s 2d Cartoons, expressive" + " "
        dtDisney50 = "" #+ "Adrian Tomine 50sAmericanCartoon" + " "
        dtKurzgesagt = "" #+ "In the style of Kurzgesagt"  + " "
        dtStoryTelling = "" #+ "Candid shot, cinematic, documentary"  + " " #For better scene composition with humans include: ||||| GOOD FOR SDXL
        dtArnonDuglas = "" #+ "In the style of Aaron Douglas" + " " #Modern Dramatic Illustrator Cartoons, shadowy             ||||| GOOD FOR SDXL
        dtAJCasson = "" #+ "In the style of A. J. Casson,  + " " #Modern Illustrator Cartoons"
        
        #### Add new variables to this list: 
        dramThingsPromptsAll = dtEksample + dtFilename +dtStyle + dtDisney80 + dtDisney50 + dtKurzgesagt + dtStoryTelling + dtArnonDuglas + dtAJCasson
        
        # Simulate image generation response
        try:
            # Add spesific text to the image generator
            if include_original_text: #set include_original_text = True, line 41.
                full_prompt = base_filename + dramThingsPromptsAll + translated_prompt_text + "\n" + unprossessed_prompt_text
                # base_filename + "" This will make shure the filename is added as the first input in the image generator
            
            # Add spesific text to the image generator, add # base_filename to the first part of image, to make theme more consistent when batch prosessing.
            else:
                full_prompt = dramThingsPromptsAll + translated_prompt_text

                
#####ADD CONSISTENT TEXT FOR IMAGE GENERATOR. #######################
#####ADD CONSISTENT TEXT FOR IMAGE GENERATOR. #######################
#####ADD CONSISTENT TEXT FOR IMAGE GENERATOR. #######################
#####ADD CONSISTENT TEXT FOR IMAGE GENERATOR. #######################

            # Generate the image using the chosen prompt
            image_paths = generate_image(full_prompt)
            print(full_prompt)
        except:
            print("Image generation server failed, skipping image generation.")
            image_paths = None
        
            # Function to brighten a color

        if image_paths:
            original_image_path, modified_image_path = image_paths
            print(f"Original Image generated at: {original_image_path}")
            print(f"Modified Image generated at: {modified_image_path}")

            avg_color = calculate_average_color(original_image_path)
            #Graied out becaouse it is not used here. Only earlyer
            brightened_color = brighten_color(avg_color, factor=3)

            # 1. Add the background image first (bottom layer)
            set_background_color(slide, original_image_path)

            # 2. Add the processed image on top of the background
            image_x = Cm(16)  # Left side of the slide
            image_y = Cm(0)   # Top of the slide
            image_width = Cm(19.05)
            image_height = Cm(19.05)

            generate_image_option = False  # Set to True to enable image generation

            if generate_image_option:
                # Original image (not blurred) on top of background
                slide.shapes.add_picture(
                    original_image_path if os.path.exists(original_image_path) else modified_image_path,
                    image_x, image_y, width=image_width, height=image_height
                )
            elif generate_image_option is None:
                # Do nothing if image generation is completely disabled
                pass
            else:
                # Add the blurred image (semi-oval blur) on top of background
                blurred_image_path = apply_semi_oval_blur(original_image_path)  # Apply feather #falloff #fade #diffuse "blur"
                slide.shapes.add_picture(
                    blurred_image_path,
                    image_x, image_y, width=image_width, height=image_height
                )

                # 3. Add the modified (background-removed) image last to be on top of all images
                slide.shapes.add_picture(
                    modified_image_path if os.path.exists(modified_image_path) else original_image_path,
                    image_x, image_y, width=image_width, height=image_height
                )

        else:
            print("No images to add to the presentation.")
            
        # Assuming `loaded_filename` is defined somewhere earlier in your code
        downloads_folder = os.path.expanduser("~/Downloads")
        base_filename = loaded_filename + " .gen"
        base_filename = re.sub(r'^.{12}', '', base_filename)  # Removes prompt_text_

        # Define folder path with the same name as the file (without extension)
        folder_path = os.path.join(downloads_folder, base_filename)

        # Create folder if it doesn't exist
        os.makedirs(folder_path, exist_ok=True)

        file_extension = ".pptx"
        counter = 1
        pptx_filename = os.path.join(folder_path, f"{base_filename}{file_extension}")

        # Ensure unique filename inside the folder
        while os.path.exists(pptx_filename):
            pptx_filename = os.path.join(folder_path, f"{base_filename}_{counter}{file_extension}")
            counter += 1

        prs.save(pptx_filename)
        print(f"Presentation saved in folder '{folder_path}' as '{pptx_filename}'")

    # Call the new function to remove lines that start with '!!!'
    remove_bulletpoints_after_exclamation(prs)

    # Set the maximum number of images you want to move on each slide
    max_images_to_move = 2  # For example, move only the first 2 images

    for slide in prs.slides:
        # Counter to track how many images have been moved
        images_moved = 0
        
        # Iterate through all shapes on the slide
        for shape in slide.shapes:
            # Check if the shape is an image (shape type 13 corresponds to Picture)
            if shape.shape_type == 13:
                # Perform actions on each image (for example, moving it behind the title)
                image = shape
                
                # Move the image behind the title (first shape) if necessary
                image._element.addprevious(slide.shapes[0]._element)
                
                # Increment the counter of moved images
                images_moved += 1
                
                # Stop moving images once the limit is reached
                if images_moved >= max_images_to_move:
                    break  # Exit the loop for this slide after moving the specified number of images
    
    # Save the PowerPoint file in the Downloads folder
    downloads_folder = os.path.expanduser("~/Downloads/" + base_filename)
    base_filename = loaded_filename + "gen "
    base_filename = re.sub(r'^.{12}', '', base_filename)  # Removes prompt_text_

    file_extension = ".pptx"
    counter = 1
    pptx_filename = os.path.join(downloads_folder, f"{base_filename}{file_extension}")

    while os.path.exists(pptx_filename):
        pptx_filename = os.path.join(downloads_folder, f"{base_filename}_{counter}{file_extension}")
        counter += 1

    # Check if the presentation has only one slide
    if len(prs.slides) == 1:
        print("Only one slide detected. Marking for reprocessing...")

        # Decrement the counter to force reprocessing the same number
        counter -= 1  # This ensures the same number is retried

        # Delete the current file since it was incomplete
      #  os.remove(pptx_filename)
      #  print(f"Incomplete presentation deleted: {pptx_filename}")

        # Return to the caller or exit the function so that the process restarts
        return  # Or use `break` if inside a loop

    # If more than one slide, finalize and save the presentation
    prs.save(pptx_filename)
    print(f"Presentation saved as {pptx_filename}")

def process_prompt(file_path, skip_llm):
    with open(file_path, 'r') as file:
        prompt_text = file.read()
    
    if skip_llm:
        
        
        print ("loaded_filename is:" + "\n" + loaded_filename)
        
        #####################SummaryForPreloadedText#######################
        #####################SummaryForPreloadedText#######################

        script_dir = os.path.dirname(os.path.abspath(__file__))  
        file_path = os.path.join(script_dir,loaded_filename)  

        if os.path.exists(file_path):
            try:
                with open(file_path, "r", encoding="utf-8") as file:
                    globalresponse_text = file.read()
            except Exception as e:
                print(f"Error reading response.txt: {e}")
        else:
            print("response.txt not found.")

        print(globalresponse_text)  # Verify content

        # Function to get a summary description from LM Studio
        def description_prompt_text(globalresponse_text):
            if not globalresponse_text.strip():
                print("No text available to summarize.")
                return "No description generated."

            request_data = {
                "messages": [
                    {"role": "system", "content": "Do not comment. Provide a descriptive one-sentence summary in english. Then one sentence with 'the purpose is likely to illustrate'"},
                    {"role": "user", "content": globalresponse_text}
                ],
                "model": "gemma-3-12b-it",
                "temperature": 0.1,
                "max_tokens": 500,
                "stream": False,
                "reset_history": True,
            }
            
            url = "http://localhost:1234/v1/chat/completions"  # Change this to your LM Studio API URL
            headers = {"Content-Type": "application/json"}

            try:
                response = requests.post(url, json=request_data, headers=headers)
                response.raise_for_status()  # Raise an error for bad responses (4xx, 5xx)
                response_json = response.json()
                return response_json.get("choices", [{}])[0].get("message", {}).get("content", "No description generated.")
            except requests.exceptions.RequestException as e:
                print(f"Error generating description: {e}")
                return "Create a response.txt file with the content of prompt_text_ to get a summary for preloaded files. Future fix: Duplicate the prompt_text_ as response.txt or better load prompt_text_XXXXX directly, have to get the name from glboal variable"

        time.sleep(1)
        twoLinesSummaryOfTheWholeTextForLMstudio = description_prompt_text(globalresponse_text)

        print("############# SUMMARY ##################" + "\n" + 
            "If gemma-3-12b-it is installed:" +  
            "Description going to LM studio to summarize the whole script: " + "\n" + "\n" +
            twoLinesSummaryOfTheWholeTextForLMstudio + "\n" + 
            "###############################" + "\n")

        #####################SummaryForpreloadedTextAbove#################

        
        
        
        
        
        
        
        
        
        
        
        
        response = prompt_text  # Load the existing text file instead
    else:
        # Define the system content in parts
        system_content_1 = "Følg disse instruksene ekstremt nøye. Ikke gi chat svar, gjennomfør kun tekstendringen som instruert. Gjør endinger også dersom teksten er lang og strukturert da dette er kravene til formattering og må være korrekt. Du lager overskrifter i en presentasjon og formaterer dem korrekt med enkel '#'"
        system_content_2 = "For hvert avsnitt lag en kort overskrift og marker den med enkel '#'"
        system_content_3 = "Bryt opp tekst som er lengre enn 50ord med # fulgt av ny linje, slik at teksten fordeles på flere sider. Ny overskirft for nye avsnitt."
        system_content_4 = "Legg til flere overskrifter, maks 300tegn pr paragraf før ny overskrift. Dersom oppgavene inkluderer spørsmål benytt føglende struktur: #Tittel, oppgave uten fremgangsmåte. #Fasit(Tittel). Ny linje, Svar, eller fremgangsmåte steg for steg dersom det er en problemløsingsoppgave. #Tittel2, oppgave2 uten fremgangsmåte. #Fasit2(Tittel2). Ny linje, Svar2, eller fremgangsmåte2 steg for steg dersom det er en problemløsingsoppgave steg for steg på ny side."
        system_content_5 = "Skriv på Norsk"

        # Combine the system content into one string
        system_content = f"{system_content_1} {system_content_2} {system_content_3} {system_content_4} {system_content_5}"

        # Build the data dictionary
        data = {
            "messages": [
                {"role": "system", "content": system_content},
                {"role": "user", "content": prompt_text}
            ],
            "config": {"gpuOffload": "max"},
            "noHup": "False"
        }
        
        response = send_request_to_localhost(data)  # Send request to LLM
        print("Response here:")
        print(response)

        # Define the filename
        base_name, ext = os.path.splitext(loaded_filename)
        save_filename = base_name + ".gen.txt"

        # Get the script's directory
        script_dir = os.path.dirname(os.path.abspath(__file__))

        save_folder = os.path.join(script_dir, "generated")  # Define the 'generated' folder path

        # Create the folder if it does not exist
        os.makedirs(save_folder, exist_ok=True)

        # Define the full save path
        save_path = os.path.join(save_folder, save_filename)

        # Save response to the file
        with open(save_path, "w", encoding="utf-8") as file:
            file.write(response)

        print(f"Response saved to {save_path}")

    if response:
        create_pptx(response)
    else:
        print("Failed to get a response from the server.")
        print("Download isnet_dis.onnx in case of error and put on goodpath")
        print("Remember to set steps: 4")
    open_draw_things()
def main():
    global loaded_filename  # To modify the global variable inside the function
    folder_path = os.path.dirname(__file__)
    #add logic for user input later
    # Ask user once whether to skip LLM processing
    user_input = input("Prosess with LM Studio? (y /no (press 'f' if already formated correctly)): ").strip().lower()
    skip_llm = user_input in ["no", "n", "f", "formatted"]
    
    # Loop through all prompt_text*.txt files
    for filename in sorted(os.listdir(folder_path)):
        if filename.startswith("prompt_text") and filename.endswith(".txt"):
            file_path = os.path.join(folder_path, filename)
            
            # Check if the file exists and process it
            if os.path.exists(file_path):
                print(f"Processing {filename}...")
                loaded_filename = filename  # Store the loaded file name in the global variable
                process_prompt(file_path, skip_llm)
            else:
                # If the file doesn't exist, create a new one with a default prompt.
                with open(file_path, 'w') as file:
                    file.write("""Write a story where Benny visits Italy Max1 Slide.
Every paragraph must have a heading #.
Every paragraph must have a body.
After every paragraph write a new line starting with "!!!" followed by detailed descriptions of characters facial expression, camera angle and scene composition for main character in one long list separated only by ","
""")
                print(f"{filename} did not exist, so a default prompt was created.")
                loaded_filename = filename  # Store the name of the newly created file

if __name__ == "__main__":
    try:
        main()
    except Exception:
        print("!!!To fix subtitle_placeholder = slide.placeholders[1] set: manuallyLoadedFirstSlide_text = (""#"" + response_text) " 
              "It is essential to force the creation of the first PowerPoint slide ")

# You can now access `loaded_filename` outside of the main function.
print(f"You have to give your .txt promptgeneration file file the suffix ""prompt_text_"" and place it in the same folder as this script")



    #Tips. Define Number of paragraphs in prompt_text_1.txt file.

#https://docs.drawthings.ai/documentation/documentation/8.scripts/
#https://lmstudio.ai/docs/api/sdk/lmstudioclient

#Legg til Transitions, med XML. Utfordringer finne rett slide, og ikke korupt fil. Vil ta litt tid.     https://stackoverflow.com/questions/73901095/python-pptx-workaround-to-add-transitions-to-slides

#subtitle_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None ##548 Skip error?