import requests  # Import requests library
import json  # Import json library
from pptx import Presentation  # Import Presentation from python-pptx
from pptx.util import Cm
from pptx.dml.color import RGBColor
import os
import base64
from PIL import Image
import numpy as np
from dis_bg_remover import remove_background  # Import dis-bg-remover
import subprocess
import time
import cv2
import traceback
import onnxruntime as ort
from pptx.util import Cm, Pt
import re
from pptx.util import Inches
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from transformers import pipeline
from PIL import Image, ImageFilter, ImageDraw
import numpy as np
import tkinter as tk
from tkinter import simpledialog

loaded_filename = None
#Må være med i følge dokumentasjonen for at ting skal bli blått
#prs = Presentation()
#bullet_slide_layout = prs.slide_layouts[1]
#slide = prs.slides.add_slide(bullet_slide_layout)
#hapes = slide.shapes
#title_shape = shapes.title
#body_shape = shapes.placeholders[1]

REMOVE_BACKGROUND = True  # Toggle for background removal

# Path to the Draw Things app
draw_things_app_path = "/Applications/Draw Things.app"
print("Remember to start API server inside of Draw Things app")

# Open Draw Things app
def open_draw_things():
    try:
        subprocess.Popen(["open", draw_things_app_path])
        print("Launching Draw Things app...")
        # Wait a few seconds to allow the app to fully initialize
        time.sleep(0.1)
    except Exception as e:
        print(f"Failed to open Draw Things app: {e}")

open_draw_things()

# Function to communicate with the local API
def send_request_to_localhost(data):
    # Check available models
    models_response = requests.get("http://localhost:1234/api/v0/models")
    if models_response.status_code == 200:
        models = models_response.json().get("data", [])
        loaded_model = next((model for model in models if model["state"] == "loaded"), None)
        
        # If no model is loaded, try to load qmistral-nemo-instruct-240
        if not loaded_model:
            print("No loaded model found. Attempting to load mistral-nemo-instruct-240...")
            model_to_load = next((model for model in models if model["id"] == "mistral-nemo-instruct-2407"), None)
            if model_to_load:
                load_response = requests.post(
                    "http://localhost:1234/api/v0/model/load",
                    json={"model": model_to_load["id"]}
                )
                if load_response.status_code == 200:
                    print("Model mistral-nemo-instruct loaded successfully.")
                    current_model = model_to_load["id"]
                else:
                    print(f"Failed to load model with status code: {load_response.status_code}")
                    return None
            else:
                print("Model qwen2.5-7b-instruct-1m not found.")
                return None
        else:
            current_model = loaded_model["id"]
    else:
        print(f"Failed to list models with status code: {models_response.status_code}")
        return None

    url = "http://localhost:1234/v1/chat/completions"
    request_data = {
        "model": current_model,
        "messages": data["messages"],
        "temperature": 0.3,
        "max_tokens": -1,
        "stream": True,
        "reset_history": True
    }
    response = requests.post(url, json=request_data, stream=True)
    response_text = ""
    if response.status_code == 200:
        try:
            for chunk in response.iter_content(chunk_size=None):
                if chunk:
                    chunk_data = chunk.decode('utf-8').strip()
                    if chunk_data.startswith("data: "):
                        chunk_data = chunk_data[6:]
                    try:
                        chunk_json = json.loads(chunk_data)
                        if "choices" in chunk_json:
                            for choice in chunk_json["choices"]:
                                if "delta" in choice and "content" in choice["delta"]:
                                    response_text += choice["delta"]["content"]
                    except json.JSONDecodeError:
                        response_text += chunk_data
                    print("Partial response:", response_text)
        except requests.exceptions.RequestException as e:
            print(f"Request exception: {e}")
        return response_text
    else:
        print(f"Request failed with status code: {response.status_code}")
        return None

# Function to generate image using Draw Things API
def generate_image(prompt):
    DRAW_THINGS_URL = 'http://127.0.0.1:7860/sdapi/v1/txt2img'
    IMG_SIZE = 512
    params = {
        "prompt": prompt,
        "negative_prompt": "(worst quality, one person, happy, low quality, normal quality, (variations):1.4), blur:1.5",
        "seed": -1,
        "steps": 25,
        "guidance_scale": 7, #Text guidance, guide, 4.5
        "batch_count": 1
    }

    
    headers = {"Content-Type": "application/json"}
    response = requests.post(DRAW_THINGS_URL, json=params, headers=headers)
    
    if response.status_code == 200:
        data = response.json()
        images = data.get("images", [])
        if images:
            temp_image_path = os.path.join("/tmp", "generated_image.png")
            with open(temp_image_path, "wb") as img_file:
                img_file.write(base64.b64decode(images[0]))

            # Attempt to remove background
            if REMOVE_BACKGROUND:
                processed_image_path = remove_background_from_image(temp_image_path)
                return (temp_image_path, processed_image_path)  # Return both paths
            
            else:
                print("Background removal disabled, using the original image.")
            
            return temp_image_path, None  # Return original and none for background removed

    else:
        print(f"Error generating image: {response.status_code}, {response.text}")

    return None, None  # Return both as None if image generation failed

def normalize(image, mean, std):
    """Normalize a numpy image with mean and standard deviation."""
    return (image / 255.0 - mean) / std

# Function to remove background from image
def remove_background_from_image(image_path):
    model_path = "/Users/sondre/Downloads/isnet_dis.onnx"
    print("Download isnet_dis.onnx in case of error")
    try:
        extracted_img, mask = remove_background(model_path, image_path)
        
        if extracted_img is None:
            print("Failed to remove background: No output image")
            return None

        extracted_img_pil = Image.fromarray(extracted_img)
        
        # Get the Downloads folder path and create a filename
        downloads_folder = os.path.expanduser("~/Downloads")
        base_filename = "processed_image_no_bg.png"
        processed_image_path = os.path.join(downloads_folder, base_filename)

        extracted_img_pil.save(processed_image_path)
        print(f"Background removed and image saved to {processed_image_path}")
        return processed_image_path
    except Exception as e:
        print(f"Failed to remove background: {e}")
        return None

import os
import cv2
import numpy as np
import onnxruntime as ort
from PIL import Image
import traceback

# Updated remove_background function
def remove_background(model_path, image_path):
    if model_path is None or image_path is None:
        return None, None

    input_size = (1024, 1024)

    try:
        # Load the ONNX model
        session = ort.InferenceSession(model_path)
        im = cv2.imread(image_path, cv2.IMREAD_COLOR)
    #    im = cv2.cvtColor(im, cv2.COLOR_BGR2RGB)  # Convert from BGR to RGB if using OpenCV

        # If image is grayscale, convert to RGB
        if len(im.shape) == 2:
            im = cv2.cvtColor(im, cv2.COLOR_GRAY2RGB)
        
        # Print initial image shape and values
        print(f"Original Image shape: {im.shape}")
        
        # Normalize the image using NumPy
        im = im.astype(np.float32)  # Convert to float
        im_normalized = normalize(im, mean=[0.5, 0.5, 0.5], std=[1.0, 1.0, 1.0])
            
        # Resize the image
        im_resized = cv2.resize(im_normalized, input_size, interpolation=cv2.INTER_LINEAR)
        im_resized = np.transpose(im_resized, (2, 0, 1))  # CHW format
        im_resized = np.expand_dims(im_resized, axis=0)  # Add batch dimension

        # Run inference
        im_resized = im_resized.astype(np.float32)  
        ort_inputs = {session.get_inputs()[0].name: im_resized}
        ort_outs = session.run(None, ort_inputs)
            
        # Process the model output
        result = ort_outs[0][0]  # Assuming single output and single batch
        result = np.clip(result, 0, 1)  # Assuming you want to clip the result to [0.5, 1] is half transparent
        result = (result * 255).astype(np.uint8)  # Rescale to [0, 255]
        result = np.transpose(result, (1, 2, 0))  # HWC format
        
        # Debug: Check result image
        print(f"Result image shape after background removal: {result.shape}")
        
        # Resize to original shape
        original_shape = im.shape[:2]
        result = cv2.resize(result, (original_shape[1], original_shape[0]), interpolation=cv2.INTER_LINEAR)

        # Ensure 'result' is 2D (H x W) and add an axis to make it (H x W x 1)
        alpha_channel = result[:, :, np.newaxis]

        # Debug: Check alpha channel
        print(f"Alpha channel shape: {alpha_channel.shape}")
        
        # Concatenate the RGB channels of 'im' with the alpha channel
        im_rgba = np.concatenate((im, alpha_channel), axis=2)

        # Debug: Check RGBA image before conversion
        print(f"RGBA Image shape: {im_rgba.shape}")
        
        # Convert to BGRA (Blue, Green, Red, Alpha)
        im_bgra = cv2.cvtColor(im_rgba, cv2.COLOR_RGBA2BGRA)

        # Debug: Check BGRA image
        print(f"BGRA Image shape: {im_bgra.shape}")
        
        # Convert to uint8 before returning
        im_bgra = im_bgra.astype(np.uint8)

        return im_bgra, result
    except Exception as e:
        print("An error occurred:")
        traceback.print_exc()
        return None, None

# Function to calculate average color
def calculate_average_color(image_path):
    try:
        image = Image.open(image_path)
        image = image.convert("RGB")
        pixels = np.array(image)
        avg_color = np.mean(pixels, axis=(0, 1))
        return tuple(avg_color.astype(int))
    except Exception as e:
        print(f"Error calculating average color: {e}")
        return (255, 255, 255)  # Default to white if there's an error

# Function to brighten a color
def brighten_color(color, factor=1.5):
    r, g, b = color
    r = min(int(r * factor), 255)
    g = min(int(g * factor), 255)
    b = min(int(b * factor), 255)
    return r, g, b

# Function to set background color for slides
def set_background_color(slide, image_path):
    try:
        avg_color = calculate_average_color(image_path)
        brightened_color = brighten_color(avg_color, factor=1.5)

        # Apply the brightened color as the slide's background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*brightened_color)
        print(f"Set background color: {brightened_color}")
    except Exception as e:
        print(f"Failed to set background color: {e}")



def apply_semi_oval_blur(image_path):
    # Open the original image
    img = Image.open(image_path)
    width, height = img.size

    # Create a new image with transparent background (RGBA mode)
    img_with_alpha = img.convert("RGBA")
    
    # Create an elliptical gradient mask with feathered edges
    mask = Image.new("L", (width, height), 255)  # Start with a white mask (full opacity)
    draw = ImageDraw.Draw(mask)

    # Loop over all pixels to create a radial gradient
    for y in range(height):
        for x in range(width):
            # Calculate the distance from the center
            dist = ((x - width / 2) ** 2 + (y - height / 2) ** 2) ** 0.5
            # Calculate the transparency based on the distance
            alpha = max(0, min(255, 255 - int(dist / (width / 2) * 255)))
            mask.putpixel((x, y), alpha)
    
    # Apply the mask to blend with alpha transparency near the edges
    img_with_alpha.putalpha(mask)  # Apply the transparency mask

    # Save the modified image
    downloads_folder = os.path.expanduser("~/Downloads")
    blurred_image_path = os.path.join(downloads_folder, "feather_falloff_image.png")

    img_with_alpha.save(blurred_image_path)  # Save the modified image

    print(f"Image with feather falloff saved to: {blurred_image_path}")

    return blurred_image_path  # Return the path to the modified image

def remove_bulletpoints_after_exclamation(prs):
    # Iterate over all slides in the presentation
    for slide in prs.slides:
        # Iterate over all shapes in the slide
        for shape in slide.shapes:
            # Check if the shape has text (e.g., a text box or placeholder)
            if hasattr(shape, "text"):
                # Split the text at the first occurrence of "!!!"
                if "!!!" in shape.text:
                    shape.text = shape.text.split("!!!")[0]
                             
            
def create_pptx(response_text):
    prs = Presentation()
    # Set the slide width and height for a 16:9 aspect ratio
    prs.slide_width = Cm(33.867)  # 16:9 width
    prs.slide_height = Cm(19.05)  # 16:9 height

    slide_layout = prs.slide_layouts[1]  # Using a predefined layout with title and subtitle placeholders
    # Handle input lines and title/subtitle extraction
    input_lines = response_text.splitlines()
    title = ""
    subtitle = ""
    skip_description = "!!!"
    
    for line in input_lines:
#        if line.startswith(skip_description):  # Use the variable in the condition
            # Skip lines that start with the value in the variable
#            print(f"Skipping line: {line}") #Uncomment to skip
#           continue  #Uncomment to skip

        if line.startswith(('#', '##', '###')):
            # If a new title is found, create a new slide
            if subtitle:
                # Add subtitle to the last slide
                #if line.startswith(skip_description):  # Use the variable in the condition
            # Skip lines that start with the value in the variable
                   # print(f"Skipping line: {line}") #Uncomment to skip
                  # continue  #Uncomment to skip
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = subtitle

                # Adjust dimensions and position for the last slide's subtitle
                subtitle_placeholder.width = Cm(17)  # Width adjustment
                subtitle_placeholder.height = Cm(15)  # Height adjustment
                subtitle_placeholder.left = Cm(1)  # Left position
                subtitle_placeholder.top = Cm(4)  # Top position

                # Set the font size of the subtitle
                text_frame = subtitle_placeholder.text_frame
                initial_font_size = 36
                max_font_size = initial_font_size
                min_font_size = 12  # Minimum font size for scaling

                subtitle_text = "".join([paragraph.text for paragraph in text_frame.paragraphs])
                text_length = len(subtitle_text)

                if text_length > 100:  # Adjust this threshold as needed
                    max_font_size = min(max_font_size - (text_length // 20), initial_font_size)

                max_font_size = max(max_font_size, min_font_size)

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(max_font_size)

                # Reset subtitle for the next title
                subtitle = ""
            
            # Create a new slide with the layout for title + subtitle
            prs.slide_width = Cm(33.87)  # 16:9 width
            prs.slide_height = Cm(19.05)  # 16:9 height
            slide = prs.slides.add_slide(slide_layout)
            title = line.lstrip('#').strip()  # Extract the title

            # Add the title to the slide
            slide.shapes.title.text = title


        else:
            # If it's not a title, accumulate it as a subtitle
            subtitle += line + "\n"

    # Add the last subtitle to the final slide after the loop
    if subtitle:
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = subtitle

        # Adjust dimensions and position for the last slide's subtitle
        subtitle_placeholder.width = Cm(18)  # Width adjustment
        subtitle_placeholder.height = Cm(15)  # Height adjustment
        subtitle_placeholder.left = Cm(1)  # Left position
        subtitle_placeholder.top = Cm(4)  # Top position

        # Set the font size of the subtitle
        text_frame = subtitle_placeholder.text_frame
        initial_font_size = 32
        max_font_size = initial_font_size
        min_font_size = 18  # Minimum font size for scaling

        subtitle_text = "".join([paragraph.text for paragraph in text_frame.paragraphs])
        text_length = len(subtitle_text)

        if text_length > 80:  # Adjust this threshold as needed
            max_font_size = min(max_font_size - (text_length // 20), initial_font_size)

        max_font_size = max(max_font_size, min_font_size)

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(max_font_size)
    
    # Now generate image for each slide based on the content (title + subtitle)
    
    

    for i, slide in enumerate(prs.slides):
        subtitle = slide.shapes.placeholders[1].text  # Get the subtitle from the placeholder
        prompt_text = "Full Body. 3D Model. Strong Emotions. Render. Cartoon for tweens. " + slide.shapes.title.text + "\n" + subtitle  # Combine title and subtitle

#####Add extra instructions for Draw Things. #######################
#####Add extra instructions for Translation.#######################
#####Add extra instructions for Scene..#######################
#####Add extra instructions for LLM..#######################
#####Translation

    def translate_prompt_text(prompt_text):
        # Send the prompt text to LM Studio for translation
        url = "http://localhost:1234/v1/chat/completions"
        request_data = {
            #LEGG TIL EKSTRA INSTRUKSJONER FOR Å OVERSETTE TIL ENGELSK
            "messages": [{"role": "system", "content": "You are given instructions to execute. Add detailed descriptions of a single anime style comic stripe panel with speachbubbles and scene composition. Give type of camera shot. Describe where details are placed within the stripe. Then add a translation of the provided text to englsih. "}],
            "model": "mistral-nemo-instruct-2407",  # Assuming the model id
            "temperature": 0.3,
            "max_tokens": 500,
            "stream": False,
            "reset_history": True,
        }
        request_data["messages"].append({"role": "user", "content": prompt_text})
        
        # Send request to LM Studio for translation
        response = requests.post(url, json=request_data)
        if response.status_code == 200:
            response_json = response.json()
            translated_text = response_json.get('choices', [{}])[0].get('message', {}).get('content', '')
            return translated_text
        else:
            print(f"Translation failed. Status code: {response.status_code}")
            return prompt_text  # Return the original prompt text if translation fails

    # Iterate through the slides to prepare and send the requests
    for i, slide in enumerate(prs.slides):
        subtitle = slide.shapes.placeholders[1].text  # Get the subtitle from the placeholder
        prompt_text = "Full Body. 3D Model. Strong Emotions. Render. Cartoon for tweens. " + slide.shapes.title.text + "\n" + subtitle  # Combine title and subtitle

        # Translate the prompt_text before sending it
        translated_prompt_text = translate_prompt_text(prompt_text)
        

        # Simulate image generation response
        try:
            image_paths = generate_image(translated_prompt_text)  # Send both title and subtitle as prompt
            print (translated_prompt_text)
        except:
            print("Image generation server failed, skipping image generation.")
            image_paths = None

        if image_paths:
            original_image_path, modified_image_path = image_paths
            print(f"Original Image generated at: {original_image_path}")
            print(f"Modified Image generated at: {modified_image_path}")

            avg_color = calculate_average_color(original_image_path)
            brightened_color = brighten_color(avg_color, factor=1.5)

            # 1. Add the background image first (bottom layer)
            set_background_color(slide, original_image_path)

            # 2. Add the processed image on top of the background
            image_x = Cm(16)  # Left side of the slide
            image_y = Cm(0)   # Top of the slide
            image_width = Cm(19.05)
            image_height = Cm(19.05)

            generate_image_option = False  # Set to True to enable image generation

            if generate_image_option:
                # Original image (not blurred) on top of background
                slide.shapes.add_picture(
                    original_image_path if os.path.exists(original_image_path) else modified_image_path,
                    image_x, image_y, width=image_width, height=image_height
                )
            elif generate_image_option is None:
                # Do nothing if image generation is completely disabled
                pass
            else:
                # Add the blurred image (semi-oval blur) on top of background
                blurred_image_path = apply_semi_oval_blur(original_image_path)  # Apply feather
                slide.shapes.add_picture(
                    blurred_image_path,
                    image_x, image_y, width=image_width, height=image_height
                )

            # 3. Add the modified (background-removed) image last to be on top of all images
            slide.shapes.add_picture(
                modified_image_path if os.path.exists(modified_image_path) else original_image_path,
                image_x, image_y, width=image_width, height=image_height
            )

        else:
            print("No images to add to the presentation.")

        # 4. Ensure text is on the topmost layer by adding it after all images
        title_shape = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]

        # Reset text and reapply to ensure it stays on the top layer
        title_shape.text = slide.shapes.title.text
        subtitle_placeholder.text = subtitle
    # Call the new function to remove lines that start with '!!!'
    remove_bulletpoints_after_exclamation(prs)

    # Save the PowerPoint file in the Downloads folder
    downloads_folder = os.path.expanduser("~/Downloads")
    base_filename = loaded_filename + "generated_PPTX"
    base_filename = re.sub(r'^.{12}', '', base_filename) #Removes prompt_text_

    file_extension = ".pptx"
    counter = 1
    pptx_filename = os.path.join(downloads_folder, f"{base_filename}{file_extension}")
    while os.path.exists(pptx_filename):
        pptx_filename = os.path.join(downloads_folder, f"{base_filename}_{counter}{file_extension}")
        counter += 1

    prs.save(pptx_filename)
    print(f"Presentation saved as {pptx_filename}")


    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Using a predefined layout with title and subtitle placeholders

    # Handle input lines and title/subtitle extraction
    input_lines = response_text.splitlines()
    title = ""
    subtitle = ""
    
    

def process_prompt(file_path):
    with open(file_path, 'r') as file:
        prompt_text = file.read()
    # Bruk !!! for å markere noe som skal sendes til LM Studio men ikke vises i presentasjon. Må være på slutten av setning.
    # Define the system content in parts
    system_content_1 = "Du lager en presentasjon."
    system_content_2 = "Hvert avsnitt inkluder: Overskrift markert med #."
    system_content_3 = "" #Structure: Max 3 Slides.
    system_content_4 = ""
    system_content_5 = "Skriv på Norsk" #!!! in the end of the paragraph

    # Combine the system content into one string
    system_content = f"{system_content_1} {system_content_2} {system_content_3} {system_content_4} {system_content_5}"

    # Define the user content
    user_content = prompt_text  # Assuming 'prompt_text' is defined earlier

    # Build the data dictionary
    data = {
        "model": "mistral-small-24b-instruct-2501",
        "messages": [
            {"role": "system", "content": system_content},
            {"role": "user", "content": user_content}
        ]
    }

    response = send_request_to_localhost(data)
    
    if response:
        create_pptx(response)
    else:
        print("Failed to get a response from the server.")
        print("Download isnet_dis.onnx in case of error and put on goodpath")
        print("Remember to set steps: 4")


def main():
    global loaded_filename  # To modify the global variable inside the functio
    folder_path = os.path.dirname(__file__)
    
    # Loop through all prompt_text*.txt files
    for filename in sorted(os.listdir(folder_path)):
        if filename.startswith("prompt_text") and filename.endswith(".txt"):
            file_path = os.path.join(folder_path, filename)
            
            # Check if the file exists and process it
            if os.path.exists(file_path):
                print(f"Processing {filename}...")
                loaded_filename = filename  # Store the loaded file name in the global variable
                process_prompt(file_path)
            else:
                # If the file doesn't exist, create a new one with a default prompt.
                with open(file_path, 'w') as file:
                    file.write("""Write a story where Benny visits Italy Max1 Slide.

Every paragraph must have a heading #.

Every paragraph must have a body.

After every paragraph write a new line starting with ""!!!"" followed by detailed descriptions of characters facial expression, camera angle and scene composition for main character in one long list separated only by "","
""")
                print(f"{filename} did not exist, so a default prompt was created.")
                loaded_filename = filename  # Store the name of the newly created file

if __name__ == "__main__":
    main()

# You can now access `loaded_filename` outside of the main function.
print(f"Last processed file: {loaded_filename}")


    #Tips. Define Number of paragraphs in prompt_text_1.txt file.
